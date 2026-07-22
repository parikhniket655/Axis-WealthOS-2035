import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def build_axis_presentation():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank slide

    # ── PALETTE ──────────────────────────────────────────────────────────────
    MAROON = RGBColor(107, 29, 46)      # Axis Bank deep burgundy (6B1D2E)
    MAROON2 = RGBColor(139, 32, 53)     # Slightly lighter burgundy (8B2035)
    GOLD = RGBColor(201, 168, 76)       # Accent gold (C9A84C)
    OFFWHITE = RGBColor(250, 246, 241)   # Warm background (FAF6F1)
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(26, 26, 46)         # Dark slate (1A1A2E)
    LIGHTGRAY = RGBColor(232, 224, 213)  # Light gray-beige (E8E0D5)
    TEXTGRAY = RGBColor(90, 74, 66)      # Charcoal body text (5A4A42)
    MINT = RGBColor(46, 125, 94)        # Green for positive metrics (2E7D5E)
    STEEL = RGBColor(52, 73, 94)        # Muted slate blue (34495E)
    DARKMAROON = RGBColor(80, 15, 30)   # Dark panel (500F1E)

    def draw_rect(slide, x, y, w, h, fill_color, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.color.rgb = fill_color  # Borderless workaround
        return shape

    def add_text_block(slide, x, y, w, h, text, fontSize, color, bold=False, italic=False, fontFace='Calibri', align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = vertical_anchor
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        lines = text.split('\n')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = fontFace
            p.font.size = Pt(fontSize)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.italic = italic
            p.alignment = align
        return txBox

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 0 — TITLE
    # ─────────────────────────────────────────────────────────────────────────
    s0 = prs.slides.add_slide(blank_layout)

    # Full dark background
    draw_rect(s0, 0, 0, 13.3, 7.5, MAROON)

    # Decorative gold diagonal band
    draw_rect(s0, 8.8, 0, 0.06, 7.5, GOLD)
    draw_rect(s0, 9.1, 0, 0.02, 7.5, GOLD)

    # Right panel darker bg
    draw_rect(s0, 8.86, 0, 4.44, 7.5, DARKMAROON)

    # AXIS BANK label top left
    add_text_block(s0, 0.5, 0.35, 4, 0.35, 'AXIS BANK', 11, GOLD, bold=True, fontFace='Calibri')

    # Main title
    add_text_block(s0, 0.5, 1.4, 7.9, 1.1, 'AXIS WEALTHOS 2035', 46, WHITE, bold=True, fontFace='Cambria')

    # Subtitle
    add_text_block(s0, 0.5, 2.65, 7.9, 1.0, 'From Relationship-Led Wealth Servicing\nto an AI-Native Household Wealth Operating System', 17, LIGHTGRAY, fontFace='Calibri')

    # Gold divider
    draw_rect(s0, 0.5, 3.75, 2.5, 0.05, GOLD)

    # Competition / date
    add_text_block(s0, 0.5, 3.9, 7.9, 0.3, 'MOVES 2026  ·  Organization of the Future: Design for 2035', 12, LIGHTGRAY, fontFace='Calibri', italic=True)

    # Team label right panel
    add_text_block(s0, 9.2, 1.6, 3.7, 0.3, 'SUBMITTED BY', 10, GOLD, bold=True, fontFace='Calibri', align=PP_ALIGN.CENTER)
    add_text_block(s0, 9.2, 2.0, 3.7, 0.5, 'Team CashCows', 22, WHITE, bold=True, fontFace='Cambria', align=PP_ALIGN.CENTER)

    # Three KPIs on right panel
    kpis = [
        {'val': '₹6.78T', 'lbl': 'Burgundy AUM\nBaseline'},
        {'val': '16,450+', 'lbl': 'Burgundy Private\nFamilies'},
        {'val': '1,700', 'lbl': 'Virtual RMs\nDeployed'},
    ]
    for i, k in enumerate(kpis):
        y = 3.1 + i * 1.3
        draw_rect(s0, 9.3, y, 3.5, 1.1, MAROON, GOLD)
        add_text_block(s0, 9.3, y + 0.08, 3.5, 0.5, k['val'], 22, GOLD, bold=True, fontFace='Cambria', align=PP_ALIGN.CENTER)
        add_text_block(s0, 9.3, y + 0.58, 3.5, 0.45, k['lbl'], 10, LIGHTGRAY, fontFace='Calibri', align=PP_ALIGN.CENTER)

    # Confidential footer
    add_text_block(s0, 0.5, 7.1, 8.2, 0.25, 'CONFIDENTIAL  ·  BOARD-LEVEL STRATEGY BRIEF  ·  JULY 2026', 9, RGBColor(158, 140, 132), fontFace='Calibri')


    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — THE BIG BET
    # ─────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    draw_rect(s1, 0, 0, 13.3, 7.5, OFFWHITE)

    # Top bar
    draw_rect(s1, 0, 0, 13.3, 0.7, MAROON)
    add_text_block(s1, 0.4, 0.1, 6, 0.45, 'SLIDE 1  ·  THE BIG BET', 10, GOLD, bold=True, fontFace='Calibri')
    add_text_block(s1, 0.4, 0.1, 12.5, 0.45, 'AXIS WEALTHOS 2035', 10, LIGHTGRAY, fontFace='Calibri', align=PP_ALIGN.RIGHT)

    # Main headline
    add_text_block(s1, 0.4, 0.82, 12.5, 0.75, 'Reinventing Wealth Management into an AI-Native Household Wealth Operating System', 21, MAROON, bold=True, fontFace='Cambria')

    # Gold underline
    draw_rect(s1, 0.4, 1.58, 12.5, 0.04, GOLD)

    # Three Columns
    colW = 3.9
    colY = 1.72
    colH = 5.35
    cols = [
        {'x': 0.35, 'color': MAROON, 'label': '01', 'head': 'THE STRUCTURAL\nCHALLENGE', 'sub': 'Why Now?'},
        {'x': 4.55, 'color': STEEL, 'label': '02', 'head': 'THE CORE BET', 'sub': 'The Transformation'},
        {'x': 8.75, 'color': MINT, 'label': '03', 'head': 'WHY AXIS WINS', 'sub': 'Strategic Anchors'},
    ]
    for c in cols:
        # Column card background
        draw_rect(s1, c['x'], colY, colW, colH, WHITE, LIGHTGRAY)
        # Colored header band
        draw_rect(s1, c['x'], colY, colW, 1.05, c['color'])
        # Number
        add_text_block(s1, c['x'] + 0.15, colY + 0.08, 0.6, 0.4, c['label'], 22, WHITE, bold=True, fontFace='Cambria')
        # Heading
        add_text_block(s1, c['x'] + 0.15, colY + 0.08, colW - 0.25, 0.55, c['head'], 13, WHITE, bold=True, fontFace='Calibri', align=PP_ALIGN.RIGHT)
        # Sub
        add_text_block(s1, c['x'] + 0.15, colY + 0.65, colW - 0.25, 0.3, c['sub'], 10, WHITE, fontFace='Calibri', italic=True, align=PP_ALIGN.RIGHT)

    # Col 1 Content — Structural Challenge
    c1x = 0.5
    c1y = colY + 1.18
    challenges = [
        {
            'head': 'Scale Impossibility',
            'body': "India's affluent class will hit 100M by 2027 (Goldman Sachs). 50%+ of new wealth creators are in Tier 2/3/4 cities. Human-only RM model cannot scale profitably outside metros.",
        },
        {
            'head': 'RM Productivity Ceiling',
            'body': 'Each RM handles 80-120 clients with manual portfolio reviews. 70% of RM time spent on prep, compliance & paperwork — not client relationships.',
        },
        {
            'head': 'Fragmented Client View',
            'body': "Axis sees only its own AUM. Client's full balance sheet — external MFs, insurance, property — is invisible. Product-pushing replaces holistic advisory.",
        },
        {
            'head': 'Cost-to-Serve Crisis',
            'body': 'High-touch RM model is economically unsustainable at scale. Relentlessly expanding headcount creates service fragmentation, not depth.',
        },
    ]
    for i, ch in enumerate(challenges):
        y = c1y + i * 1.03
        draw_rect(s1, c1x, y, colW - 0.3, 0.95, RGBColor(253, 240, 240), RGBColor(232, 208, 208))
        add_text_block(s1, c1x + 0.1, y + 0.05, colW - 0.45, 0.25, ch['head'], 10, MAROON, bold=True, fontFace='Calibri')
        add_text_block(s1, c1x + 0.1, y + 0.28, colW - 0.45, 0.62, ch['body'], 8.5, TEXTGRAY, fontFace='Calibri')

    # Col 2 Content — The Core Bet
    c2x = 4.7
    c2y = colY + 1.15
    # Tagline band
    draw_rect(s1, c2x, c2y, colW - 0.3, 0.45, STEEL)
    add_text_block(s1, c2x + 0.1, c2y + 0.05, colW - 0.45, 0.36, 'AXIS WEALTHOS: Stop managing accounts. Start orchestrating financial lives.', 8.5, WHITE, bold=True, fontFace='Calibri', italic=True)

    layers = [
        {'lbl': 'DATA LAYER', 'head': 'Wealth Twin', 'body': 'Encrypted real-time financial graph — internal bank data + external consented holdings via Account Aggregator (AA) APIs. Maps full household balance sheet across generations.'},
        {'lbl': 'INTELLIGENCE', 'head': 'AI Advice Factory', 'body': 'Autonomous multi-agent backend continuously scanning the Wealth Twin. Flags portfolio drift, tax-harvesting windows, ALM mismatches, concentrated stock risk, and succession triggers — in real time.'},
        {'lbl': 'DELIVERY', 'head': 'Pod-Based Advisory', 'body': 'Replaces isolated RM model. Each client served by a Lead Relationship Partner + Portfolio Architect + Credit Specialist. Pods operate on AI-generated briefing dossiers — prep in minutes, not hours.'},
        {'lbl': 'ECOSYSTEM', 'head': 'One Axis Orchestration', 'body': 'WealthOS routes seamlessly to Axis Securities, AMC, Capital, Finance & Trustee. A client liquidity event triggers coordinated action across all entities — not siloed product conversations.'},
    ]
    for i, l in enumerate(layers):
        y = c2y + 0.55 + i * 1.18
        draw_rect(s1, c2x, y, colW - 0.3, 1.1, RGBColor(240, 244, 248), RGBColor(200, 216, 232))
        draw_rect(s1, c2x, y, 0.06, 1.1, STEEL)
        add_text_block(s1, c2x + 0.15, y + 0.05, colW - 0.45, 0.2, l['lbl'], 7.5, STEEL, bold=True, fontFace='Calibri')
        add_text_block(s1, c2x + 0.15, y + 0.23, colW - 0.45, 0.22, l['head'], 11, DARK, bold=True, fontFace='Cambria')
        add_text_block(s1, c2x + 0.15, y + 0.44, colW - 0.45, 0.6, l['body'], 8, TEXTGRAY, fontFace='Calibri')

    # Col 3 Content — Why Axis Wins
    c3x = 8.9
    c3y = colY + 1.15
    # Big stat callout
    draw_rect(s1, c3x, c3y, colW - 0.3, 1.1, MINT)
    add_text_block(s1, c3x, c3y + 0.05, colW - 0.3, 0.55, '₹6.78T', 36, WHITE, bold=True, fontFace='Cambria', align=PP_ALIGN.CENTER)
    add_text_block(s1, c3x, c3y + 0.65, colW - 0.3, 0.35, 'Burgundy AUM  ·  20% YoY Growth  ·  FY26', 9, WHITE, fontFace='Calibri', align=PP_ALIGN.CENTER)

    anchors = [
        {'head': 'Citibank Integration', 'body': "₹1.11T wealth AUM added. Premium talent pool + disciplined UHNI base instantly elevates Axis's premiumization strategy."},
        {'head': 'Virtual RM Infrastructure', 'body': '1,700 virtual RMs serving 4.2M clients/month. Proves Axis already has the digital service backbone to scale automated advisory.'},
        {'head': 'One Axis Multi-Entity Edge', 'body': 'Securities + AMC + Capital + Finance + Trustee. No standalone wealth boutique or competitor can replicate this integrated group capability.'},
        {'head': 'AA Regulatory Tailwind', 'body': "RBI's Account Aggregator framework provides legal, real-time consent-based data rails. Axis can legally view a client's FULL balance sheet — including competitor holdings."},
    ]
    for i, a in enumerate(anchors):
        y = c3y + 1.2 + i * 1.0
        draw_rect(s1, c3x, y, colW - 0.3, 0.93, RGBColor(239, 248, 244), RGBColor(184, 222, 206))
        draw_rect(s1, c3x, y, 0.06, 0.93, MINT)
        add_text_block(s1, c3x + 0.15, y + 0.05, colW - 0.45, 0.22, a['head'], 10, MINT, bold=True, fontFace='Calibri')
        add_text_block(s1, c3x + 0.15, y + 0.27, colW - 0.45, 0.6, a['body'], 8.5, TEXTGRAY, fontFace='Calibri')

    # Speaker note bottom banner
    draw_rect(s1, 0.35, 7.08, 12.6, 0.32, MAROON)
    add_text_block(s1, 0.5, 7.09, 12.3, 0.28, '"Burgundy is already a market leader. But to capture the next wave of India\'s wealth, we must multiply advisor capability — not headcount. WealthOS is that future."', 9, GOLD, italic=True, align=PP_ALIGN.CENTER)


    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — HOW IT WORKS & IMPACT
    # ─────────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    draw_rect(s2, 0, 0, 13.3, 7.5, OFFWHITE)

    # Top bar
    draw_rect(s2, 0, 0, 13.3, 0.7, MAROON)
    add_text_block(s2, 0.4, 0.1, 8, 0.45, 'SLIDE 2  ·  HOW IT WORKS & STRATEGIC IMPACT', 10, GOLD, bold=True, fontFace='Calibri')
    add_text_block(s2, 0.4, 0.1, 12.5, 0.45, 'AXIS WEALTHOS 2035', 10, LIGHTGRAY, fontFace='Calibri', align=PP_ALIGN.RIGHT)

    # Main headline
    add_text_block(s2, 0.4, 0.82, 12.5, 0.55, 'Advisory Pods · Multi-Entity Integration · Business Case', 20, MAROON, bold=True, fontFace='Cambria')
    draw_rect(s2, 0.4, 1.38, 12.5, 0.04, GOLD)

    # Left Column Flowchart Setup
    LX = 0.35
    LY = 1.5
    LW = 6.7
    add_text_block(s2, LX, LY, LW, 0.28, 'THE WEALTHOS OPERATING MODEL', 10, MAROON, bold=True, fontFace='Calibri')

    # Flow Steps
    flowSteps = [
        {
            'label': 'INPUT',
            'color': STEEL,
            'head': 'Client Financial Data',
            'items': ['Internal: Axis deposits, credit, card spends', 'External (AA APIs): MFs, equities, insurance,\nliabilities, property, business cash flows'],
        },
        {
            'label': 'INTELLIGENCE',
            'color': STEEL,
            'head': 'AI Advice Factory',
            'items': ['Continuous event detection: drift, tax-harvest,\nALM mismatch, succession triggers', 'Risk suitability audit + explainability log\nbefore any recommendation surfaces'],
        },
        {
            'label': 'DELIVERY',
            'color': MAROON2,
            'head': 'Advisory Pod',
            'items': ['Lead RM (trust + goal discovery)', 'Portfolio Architect (asset allocation)', 'Credit Specialist (leverage + liquidity)'],
        },
        {
            'label': 'EXECUTION',
            'color': MINT,
            'head': 'One Axis Ecosystem',
            'items': ['Axis Securities  ·  Axis AMC  ·  Axis Capital', 'Axis Finance  ·  Axis Trustee Services', 'Seamless group routing — single client view'],
        },
    ]

    fBoxW = (LW / 4) - 0.12
    for i, f in enumerate(flowSteps):
        fx = LX + i * (fBoxW + 0.12)
        fy = LY + 0.38
        fh = 2.55

        # Draw box
        draw_rect(s2, fx, fy, fBoxW, fh, WHITE, LIGHTGRAY)
        # Header band
        draw_rect(s2, fx, fy, fBoxW, 0.48, f['color'])
        add_text_block(s2, fx + 0.08, fy + 0.04, fBoxW - 0.1, 0.18, f['label'], 7.5, WHITE, bold=True, fontFace='Calibri')
        add_text_block(s2, fx + 0.08, fy + 0.22, fBoxW - 0.1, 0.22, f['head'], 9.5, WHITE, bold=True, fontFace='Cambria')

        for j, item in enumerate(f['items']):
            add_text_block(s2, fx + 0.08, fy + 0.56 + j * 0.65, fBoxW - 0.12, 0.6, '▸  ' + item, 8, TEXTGRAY, fontFace='Calibri')

        # Arrow between boxes
        if i < 3:
            draw_rect(s2, fx + fBoxW, fy + (fh / 2) - 0.04, 0.12, 0.08, GOLD)

    # Phase rollout roadmap
    draw_rect(s2, LX, LY + 3.07, LW, 0.26, MAROON)
    add_text_block(s2, LX + 0.15, LY + 3.1, LW - 0.2, 0.2, 'IMPLEMENTATION ROADMAP', 9, GOLD, bold=True, fontFace='Calibri')

    phases = [
        {'ph': 'PHASE 1  ·  Yr 1–2', 'col': MAROON2, 'head': 'Foundation & HNI Pilot', 'items': ['Wealth Twin data layer build', 'AA API integration', '500 RM pilot deployment', 'Advisory productivity lift proof']},
        {'ph': 'PHASE 2  ·  Yr 3–5', 'col': STEEL, 'head': 'Affluent Scale-Up', 'items': ['Expand to full Burgundy segment', 'Family Banking 2.0 launch', 'Pod-based coverage rollout', 'GST + business cash flow data']},
        {'ph': 'PHASE 3  ·  Yr 6–10', 'col': MINT, 'head': 'One Axis Integration', 'items': ['Full entity routing activated', 'Open API ecosystem partners', 'Cross-border & AIF integration', '100% data-native household advisory']},
    ]
    phW = (LW / 3) - 0.1
    for i, p in enumerate(phases):
        px = LX + i * (phW + 0.1)
        py = LY + 3.4
        draw_rect(s2, px, py, phW, 2.55, WHITE, LIGHTGRAY)
        draw_rect(s2, px, py, phW, 0.38, p['col'])
        add_text_block(s2, px + 0.08, py + 0.04, phW - 0.12, 0.2, p['ph'], 7.5, WHITE, bold=True, fontFace='Calibri')
        add_text_block(s2, px + 0.08, py + 0.4, phW - 0.12, 0.28, p['head'], 10, DARK, bold=True, fontFace='Cambria')
        for j, item in enumerate(p['items']):
            add_text_block(s2, px + 0.1, py + 0.72 + j * 0.44, phW - 0.15, 0.4, '•  ' + item, 8, TEXTGRAY, fontFace='Calibri')

    # Right Column Setup (Impact Table + Feasibility + Risks)
    RX = 7.2
    RY = 1.5
    RW = 5.75
    add_text_block(s2, RX, RY, RW, 0.28, 'FINANCIAL & FEASIBILITY SCORECARD', 10, MAROON, bold=True, fontFace='Calibri')

    # Table headers
    tY = RY + 0.32
    draw_rect(s2, RX, tY, RW, 0.3, MAROON)
    hCols = ['INDICATOR', 'FY26 BASE', '2030', '2035']
    hWids = [2.0, 1.15, 1.15, 1.15]
    hOff = RX + 0.08
    for i, h in enumerate(hCols):
        add_text_block(s2, hOff, tY + 0.05, hWids[i], 0.2, h, 8, GOLD, bold=True, fontFace='Calibri')
        hOff += hWids[i] + 0.05

    rows = [
        ['Burgundy AUM', '₹6.78T', '₹11.50T', '₹16.50T'],
        ['Burgundy Private AUM', '₹2.40T', '₹4.50T', '₹7.00T'],
        ['UHNWI Families Served', '16,450+', '30,000+', '55,000+'],
        ['RM Productivity', '1.0x', '1.3x', '1.6x'],
        ['Wallet Share (Affluent)', '~18%', '25%', '35%'],
        ['Cost-to-Serve', 'Baseline', '–35%', '–85%'],
        ['Advisory Turnaround', 'T+2 days', 'T+4 hrs', '<5 min'],
    ]

    for i, r in enumerate(rows):
        ry = tY + 0.3 + i * 0.37
        bg = RGBColor(245, 240, 235) if i % 2 == 0 else WHITE
        draw_rect(s2, RX, ry, RW, 0.37, bg)
        xOff = RX + 0.08
        for j, cell in enumerate(r):
            isPositive = (j >= 2) and any(x in cell for x in ['T', '%', 'min', 'hrs', '+', 'x'])
            isNegative = cell.startswith('–')
            
            cell_color = DARK
            if isPositive or isNegative:
                cell_color = MINT
            elif j == 0:
                cell_color = TEXTGRAY
                
            add_text_block(s2, xOff, ry + 0.07, hWids[j], 0.24, cell, 
                           8.5 if j == 0 else 9, cell_color, bold=(j > 0), fontFace='Calibri')
            xOff += hWids[j] + 0.05

    # Feasibility Scoreboard
    feasY = tY + 0.3 + len(rows) * 0.37 + 0.15
    draw_rect(s2, RX, feasY, RW, 0.28, STEEL)
    add_text_block(s2, RX + 0.1, feasY + 0.05, RW - 0.15, 0.18, 'FEASIBILITY ASSESSMENT', 9, WHITE, bold=True, fontFace='Calibri')

    feases = [
        {'lbl': 'Strategic Fit', 'score': '9/10'},
        {'lbl': 'Market Demand', 'score': '9/10'},
        {'lbl': 'Tech & Data', 'score': '8/10'},
        {'lbl': 'Operational', 'score': '7/10'},
        {'lbl': 'Regulatory', 'score': '8/10'},
    ]
    fFW = RW / len(feases)
    for i, f in enumerate(feases):
        fx = RX + i * fFW
        fy = feasY + 0.3
        draw_rect(s2, fx + 0.04, fy, fFW - 0.08, 0.7, WHITE, LIGHTGRAY)
        add_text_block(s2, fx + 0.04, fy + 0.06, fFW - 0.08, 0.3, f['score'], 16, MAROON, bold=True, fontFace='Cambria', align=PP_ALIGN.CENTER)
        add_text_block(s2, fx + 0.04, fy + 0.38, fFW - 0.08, 0.26, f['lbl'], 7.5, TEXTGRAY, fontFace='Calibri', align=PP_ALIGN.CENTER)

    # Risk Mitigations
    riskY = feasY + 1.1
    draw_rect(s2, RX, riskY, RW, 0.28, MAROON2)
    add_text_block(s2, RX + 0.1, riskY + 0.05, RW - 0.15, 0.18, 'KEY RISK MITIGATIONS', 9, WHITE, bold=True, fontFace='Calibri')

    risks = [
        {'risk': 'RM Resistance', 'mit': 'AI as co-pilot: automates prep & compliance; RMs own 100% relationship time'},
        {'risk': 'Data Privacy', 'mit': 'Granular client-revocable consent; zero-knowledge proofs; live data dashboards'},
        {'risk': 'Model Bias / Hallucination', 'mit': 'Human-in-the-loop: Portfolio Architect validates every AI proposal before client'},
        {'risk': 'One Axis Friction', 'mit': 'Unified Group Wallet Share scorecard; shared commissions across all entities'},
    ]
    for i, r in enumerate(risks):
        ry = riskY + 0.32 + i * 0.52
        bg = RGBColor(253, 240, 240) if i % 2 == 0 else WHITE
        draw_rect(s2, RX, ry, RW, 0.5, bg, LIGHTGRAY)
        add_text_block(s2, RX + 0.1, ry + 0.04, 1.5, 0.42, '⚠ ' + r['risk'] + ':', 8, MAROON, bold=True, fontFace='Calibri')
        add_text_block(s2, RX + 1.6, ry + 0.1, RW - 1.65, 0.35, r['mit'], 8, TEXTGRAY, fontFace='Calibri')

    # Footer banner
    draw_rect(s2, 0, 7.17, 13.3, 0.33, MAROON)
    add_text_block(s2, 0.4, 7.18, 12.5, 0.28, '"WealthOS is not an app or a chatbot. It is a fundamental operational rewiring — turning Axis\'s multi-entity capabilities into one cohesive client experience."', 9, GOLD, italic=True, align=PP_ALIGN.CENTER)

    # Save
    out_path = '/Users/ketanparikh/Desktop/Axis/AxisWealthOS2035_TeamCashCows.pptx'
    prs.save(out_path)
    print(f"Presentation saved successfully to: {out_path}")

if __name__ == "__main__":
    build_axis_presentation()
